from dvc_mapping import busMark, branchMark
from plot_utils import get_buscoords, get_busdetails, reformat
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Mm, Pt, Inches
from pptx.dml.color import RGBColor

from psse_utils import run_psse
businfo, mybusdat, rgenbus, myloadinfo, brnflow, trfflow = run_psse()

# Create an empty presentation with no slides

# Add a slide with the slide layout #6.
# Some other layouts are:
#   0 - Title slide
#   1 - Title and Content
#   2 - Section Header
# ... etc (in PowerPoint, click the New Slide dropdown to see the layouts)


# Powerflow plot sized for A3 paper (420x297)mm ~ (16.5x11.7)in ~ (1188x842)pt
# Slide default size is 10"x7.5" (254x190)mm ~ (1024x768)pixels @ 96dpi
# Viewport/slide coordinates : (720x540)pt (w x h)
shp_width = Mm(3.5)
shp_height = Mm(3.5)
txt_width = Pt(1)
txt_height = Pt(1)

def pptx_bus_trace(bdat):
	for bus in bdat:
		if bus in busMark:
			dot, busname = get_busdetails(businfo,bus)
			x, y = get_buscoords(bus)
			shp_left = Mm((254./420.)*x)-(shp_width/2.)
			shp_top = Mm(190.-(190./297.)*y)-(shp_height/2.)
			oshp = shapes.add_shape( MSO_SHAPE.OVAL, shp_left, shp_top, 
				shp_width, shp_height)
			colr, busname = get_busdetails(businfo,bus)
			r,g,b = reformat(colr)
			fill = oshp.fill
			fill.solid()
			fill.fore_color.rgb = RGBColor(r,g,b)
			txt_left = shp_left+Pt(9)
			txt_top = shp_top
			txBox = shapes.add_textbox(txt_left, txt_top, txt_width, txt_height)
			tf = txBox.text_frame
			tf.clear()
			p = tf.paragraphs[0]
			run = p.add_run()
			run.text = busname
			font = run.font
			font.name = 'Calibri'
			font.size = Pt(8)

def pptx_brn_trace():
    for brn in brnflow:
        if brn[0] in busMark and brn[1] in busMark:
            key = (brn[0],brn[1])
            rkey = tuple(reversed(key))
            colr, busname = get_busdetails(businfo,brn[0])
            place_lines(key,'brn',colr)
            if brn[4]>0: # check powerflow direction
                args = (key,0.5,6,colr,brn[2])
            else:
                args = (rkey,0.5,6,colr,brn[2])
            place_arrows_pfdata(*args)

        
if __name__ == '__main__':
	prs = Presentation()
	title_slide_layout = prs.slide_layouts[0]
	slide = prs.slides.add_slide(title_slide_layout)
	title = slide.shapes.title
	subtitle = slide.placeholders[1]
	title.text = "Damodar Valley Corporation"
	subtitle.text = "System Planning & Engineering"
	slide = prs.slides.add_slide(prs.slide_layouts[6])
	shapes = slide.shapes
	pptx_bus_trace(mybusdat)  
    #trn_trace(trfflow)
    #bus_trace(mybusdat)
	prs.save('dvc.pptx')

